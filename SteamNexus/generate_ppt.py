import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def create_presentation():
    # 1. Initialize Presentation
    prs = Presentation()
    
    # helper for adding slides
    def add_slide(title_text, content_points=None, image_path=None):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Set Title
        title = slide.shapes.title
        title.text = title_text
        
        # Set Body
        if content_points:
            tf = slide.placeholders[1].text_frame
            tf.text = content_points[0]
            for point in content_points[1:]:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0

    # Slide 1: Title
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "SteamNexus"
    subtitle.text = "Inteligencia de Datos para el Descubrimiento de Videojuegos\nBig Data Project - 2026"

    # Slide 2: El Problema
    add_slide("¿Por qué SteamNexus?", [
        "Parálisis por Análisis: +100,000 juegos en la plataforma.",
        "Sesgo de Popularidad: Los algoritmos actuales ignoran joyas independientes.",
        "Dificultad para encontrar nichos específicos basados en mecánicas complejas."
    ])

    # Slide 3: La Solución (Los dos datas)
    add_slide("Ecosistema de Datos (The Two Datasets)", [
        "Dataset 1: Metadatos de Juegos (Games Metadata) - 80k+ títulos.",
        "Dataset 2: Reseñas de Usuarios (6.4M Reviews) - Sentimiento real.",
        "Sinergia: Unimos la ficha técnica con la voz de la comunidad."
    ])

    # Slide 4: Arquitectura Técnica
    add_slide("Arquitectura del Proyecto", [
        "Ingesta: Automatizada vía Kaggle API.",
        "Limpieza: Normalización de campos y corrección de desplazamientos.",
        "Análisis: Grafos de tags y NLP para análisis de sentimientos.",
        "Visualización: Dashboards interactivos y grafos de redes."
    ])

    # Save
    output_name = "SteamNexus_Presentation.pptx"
    prs.save(output_name)
    print(f"OK: Presentation created successfully: {output_name}")

if __name__ == "__main__":
    create_presentation()
